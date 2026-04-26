"""Content extraction for Office documents and text files.

Primary: Microsoft markitdown (DOCX, PPTX, XLSX, PDF, MSG → Markdown).
Fallback: python-docx, openpyxl, python-pptx for specific formats.
HTML → text conversion for OneNote and email bodies.
Optional OCR fallback (Tesseract) for scanned PDFs with no text layer.

Top-level imports are stdlib-only; all heavy dependencies are lazy-loaded.
"""

from __future__ import annotations

import logging
import os
import tempfile
from pathlib import Path

log = logging.getLogger("pb-content")

# Office document extensions that markitdown can handle
MARKITDOWN_EXTENSIONS = frozenset({
    ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls",
    ".pdf", ".msg", ".eml", ".rtf",
})

# Binary extensions to skip entirely (images, video, archives)
BINARY_SKIP = frozenset({
    ".png", ".jpg", ".jpeg", ".gif", ".ico", ".svg", ".webp", ".bmp",
    ".mp3", ".mp4", ".wav", ".avi", ".mov", ".mkv", ".wmv",
    ".zip", ".tar", ".gz", ".bz2", ".xz", ".7z", ".rar",
    ".exe", ".dll", ".so", ".dylib", ".bin", ".msi",
    ".iso", ".dmg", ".img",
    ".woff", ".woff2", ".ttf", ".eot", ".otf",
    ".pyc", ".pyo", ".class", ".o", ".a",
    ".db", ".sqlite", ".sqlite3",
})

# Text-based extensions that can be read directly
TEXT_EXTENSIONS = frozenset({
    ".md", ".markdown", ".rst", ".txt", ".csv", ".tsv",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".go", ".rs", ".rb",
    ".c", ".cpp", ".h", ".hpp", ".cs", ".php", ".swift", ".kt", ".scala",
    ".sh", ".bash", ".zsh", ".ps1",
    ".yaml", ".yml", ".json", ".toml", ".xml", ".html", ".css", ".scss",
    ".sql", ".graphql", ".proto", ".rego",
    ".dockerfile", ".tf", ".hcl",
    ".env", ".ini", ".cfg", ".conf",
})

# Content type mapping for Powerbrain
CONTENT_TYPE_MAP: dict[str, str] = {
    ".md": "markdown", ".markdown": "markdown",
    ".docx": "markdown", ".doc": "markdown",  # markitdown converts to markdown
    ".pptx": "markdown", ".ppt": "markdown",
    ".xlsx": "markdown", ".xls": "markdown",
    ".pdf": "markdown",
    ".msg": "markdown", ".eml": "markdown",
    ".html": "html", ".htm": "html",
    ".txt": "text", ".csv": "csv", ".tsv": "tsv",
    ".json": "json", ".yaml": "yaml", ".yml": "yaml",
    ".xml": "xml",
}

# ── OCR fallback configuration ──
# Default-off because Tesseract and pdf2image add ~120 MB + poppler system deps.
_OCR_FALLBACK_ENABLED = os.getenv("OCR_FALLBACK_ENABLED", "false").lower() in ("1", "true", "yes")
_OCR_FALLBACK_MIN_CHARS = int(os.getenv("OCR_FALLBACK_MIN_CHARS", "50"))
_OCR_FALLBACK_MAX_PAGES = int(os.getenv("OCR_FALLBACK_MAX_PAGES", "20"))
_OCR_FALLBACK_DPI = int(os.getenv("OCR_FALLBACK_DPI", "200"))


def detect_content_type(filename: str) -> str:
    """Detect content type from filename."""
    _, ext = os.path.splitext(filename.lower())
    return CONTENT_TYPE_MAP.get(ext, "text")


def should_skip_file(filename: str) -> bool:
    """Check if a file should be skipped based on extension."""
    _, ext = os.path.splitext(filename.lower())
    return ext in BINARY_SKIP


def can_extract(filename: str) -> bool:
    """Check if we can extract text from this file type."""
    _, ext = os.path.splitext(filename.lower())
    return ext in MARKITDOWN_EXTENSIONS or ext in TEXT_EXTENSIONS


class ContentExtractor:
    """Extract text content from Office documents and other file types.

    Usage:
        extractor = ContentExtractor()
        text = extractor.extract_from_bytes(data, "report.pdf")

    The `extract_from_bytes_detailed` variant additionally reports which backend
    produced the text (markitdown / fallback / text / ocr) — useful for telemetry.
    """

    def __init__(self):
        self._markitdown = None

    # ── Public API ───────────────────────────────────────────────

    def extract_from_bytes(self, data: bytes, filename: str) -> str | None:
        """Extract text content from file bytes.

        Returns markdown/text string, or None if extraction fails or file should be skipped.
        """
        text, _ = self.extract_from_bytes_detailed(data, filename)
        return text

    def extract_from_bytes_detailed(
        self, data: bytes, filename: str
    ) -> tuple[str | None, str]:
        """Like extract_from_bytes but also returns the backend used.

        Returns a tuple (text, extractor) where extractor is one of:
          "markitdown" | "fallback" | "text" | "ocr" | "skipped" | "failed"
        """
        if should_skip_file(filename):
            return None, "skipped"

        _, ext = os.path.splitext(filename.lower())

        # Text files: decode directly
        if ext in TEXT_EXTENSIONS:
            try:
                return data.decode("utf-8"), "text"
            except UnicodeDecodeError:
                log.debug("Cannot decode %s as UTF-8, skipping", filename)
                return None, "failed"

        # Office documents: use markitdown
        if ext in MARKITDOWN_EXTENSIONS:
            text, extractor = self._extract_with_markitdown_detailed(data, filename)
            # OCR fallback for PDFs whose extracted text is suspiciously short
            if (
                ext == ".pdf"
                and _OCR_FALLBACK_ENABLED
                and (not text or len(text.strip()) < _OCR_FALLBACK_MIN_CHARS)
            ):
                ocr_text = self._extract_pdf_ocr(data, filename)
                if ocr_text and len(ocr_text.strip()) >= _OCR_FALLBACK_MIN_CHARS:
                    return ocr_text, "ocr"
            return text, extractor

        # Unknown extension: try UTF-8 decode
        try:
            return data.decode("utf-8"), "text"
        except UnicodeDecodeError:
            return None, "failed"

    def extract_html_to_text(self, html: str) -> str:
        """Convert HTML to plain text. Used for OneNote pages and email bodies."""
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(html, "html.parser")

            # Remove script and style elements
            for element in soup(["script", "style"]):
                element.decompose()

            text = soup.get_text(separator="\n")
            # Collapse multiple blank lines
            lines = [line.strip() for line in text.splitlines()]
            text = "\n".join(line for line in lines if line)
            return text
        except ImportError:
            log.warning("beautifulsoup4 not installed, using regex HTML stripping")
            import re as _re
            # Remove script/style blocks (including content), then strip tags
            text = _re.sub(r"<script\b[^>]*>.*?</script>", "", html, flags=_re.IGNORECASE | _re.DOTALL)
            text = _re.sub(r"<style\b[^>]*>.*?</style>", "", text, flags=_re.IGNORECASE | _re.DOTALL)
            text = _re.sub(r"<[^>]+>", " ", text)
            return " ".join(text.split())

    # ── Internal: markitdown ─────────────────────────────────────

    def _get_markitdown(self):
        """Lazy-init markitdown converter."""
        if self._markitdown is None:
            try:
                from markitdown import MarkItDown
                self._markitdown = MarkItDown()
            except ImportError:
                log.warning(
                    "markitdown not installed. Office document extraction disabled. "
                    "Install with: pip install markitdown"
                )
                raise
        return self._markitdown

    def _extract_with_markitdown_detailed(
        self, data: bytes, filename: str
    ) -> tuple[str | None, str]:
        """Extract text using markitdown; return (text, extractor_name)."""
        try:
            converter = self._get_markitdown()
        except ImportError:
            text = self._extract_fallback(data, filename)
            return text, ("fallback" if text else "failed")

        suffix = Path(filename).suffix
        tmp_path: str | None = None
        try:
            with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
                tmp.write(data)
                tmp_path = tmp.name

            result = converter.convert(tmp_path)
            text = result.text_content
            if text and text.strip():
                return text.strip(), "markitdown"
            return None, "failed"
        except Exception:
            log.warning("markitdown failed for %s, trying fallback", filename, exc_info=True)
            text = self._extract_fallback(data, filename)
            return text, ("fallback" if text else "failed")
        finally:
            if tmp_path is not None:
                try:
                    os.unlink(tmp_path)
                except OSError:
                    pass

    # ── Internal: per-format fallbacks ───────────────────────────

    def _extract_fallback(self, data: bytes, filename: str) -> str | None:
        """Fallback extraction using python-docx, openpyxl, python-pptx."""
        _, ext = os.path.splitext(filename.lower())

        try:
            if ext in (".docx", ".doc"):
                return self._extract_docx(data)
            elif ext in (".xlsx", ".xls"):
                return self._extract_xlsx(data)
            elif ext in (".pptx", ".ppt"):
                return self._extract_pptx(data)
        except ImportError:
            log.warning("Fallback library not installed for %s", ext)
        except Exception:
            log.warning("Fallback extraction failed for %s", filename, exc_info=True)

        return None

    @staticmethod
    def _extract_docx(data: bytes) -> str | None:
        """Extract text from DOCX using python-docx."""
        import io
        from docx import Document

        doc = Document(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        if not paragraphs:
            return None
        return "\n\n".join(paragraphs)

    @staticmethod
    def _extract_xlsx(data: bytes) -> str | None:
        """Extract text from XLSX as markdown tables using openpyxl."""
        import io
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts: list[str] = []

        for sheet in wb.worksheets:
            rows = list(sheet.iter_rows(values_only=True))
            if not rows:
                continue

            parts.append(f"## {sheet.title}\n")
            for row in rows:
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    parts.append("| " + " | ".join(cells) + " |")

        wb.close()
        text = "\n".join(parts)
        return text if text.strip() else None

    @staticmethod
    def _extract_pptx(data: bytes) -> str | None:
        """Extract text from PPTX using python-pptx."""
        import io
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        parts: list[str] = []

        for i, slide in enumerate(prs.slides, 1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                parts.append(f"## Slide {i}\n\n" + "\n\n".join(slide_texts))

        text = "\n\n".join(parts)
        return text if text.strip() else None

    # ── Internal: OCR fallback ───────────────────────────────────

    @staticmethod
    def _extract_pdf_ocr(data: bytes, filename: str) -> str | None:
        """OCR fallback for scanned PDFs using Tesseract.

        Only active when OCR_FALLBACK_ENABLED=true and the corresponding
        system packages (tesseract-ocr, poppler-utils) are installed.
        """
        try:
            import pytesseract
            from pdf2image import convert_from_bytes
        except ImportError:
            log.warning(
                "OCR fallback requested but pytesseract/pdf2image not installed for %s",
                filename,
            )
            return None

        try:
            images = convert_from_bytes(
                data,
                dpi=_OCR_FALLBACK_DPI,
                last_page=_OCR_FALLBACK_MAX_PAGES,
            )
        except Exception:
            log.warning("OCR: pdf2image failed for %s", filename, exc_info=True)
            return None

        if not images:
            return None

        texts: list[str] = []
        for idx, image in enumerate(images, 1):
            try:
                page_text = pytesseract.image_to_string(image) or ""
            except Exception:
                log.warning(
                    "OCR: tesseract failed on page %d of %s", idx, filename, exc_info=True
                )
                continue
            if page_text.strip():
                texts.append(f"## Page {idx}\n\n{page_text.strip()}")

        if not texts:
            return None
        return "\n\n".join(texts)
